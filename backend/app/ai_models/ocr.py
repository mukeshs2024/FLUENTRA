from pathlib import Path

from fastapi import UploadFile

from app.utils.tempfiles import save_upload_to_temp


def extract_text(file: UploadFile) -> str:
    temp_path = save_upload_to_temp(file)
    try:
        ext = temp_path.suffix.lower()
        if ext == ".pdf":
            return _extract_pdf_text(temp_path)
        if ext == ".docx":
            return _extract_docx_text(temp_path)
        if ext in {".ppt", ".pptx"}:
            return _extract_pptx_text(temp_path)
        if ext in {".txt"}:
            return temp_path.read_text(encoding="utf-8", errors="ignore")
        if ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff"}:
            return _extract_image_text(temp_path)
        return ""
    finally:
        temp_path.unlink(missing_ok=True)


def _extract_pdf_text(path: Path) -> str:
    import fitz
    import pdfplumber
    import pytesseract
    from PIL import Image

    text_blocks = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text_blocks.append(page_text)
    combined = "\n".join(text_blocks).strip()
    if combined:
        return combined

    doc = fitz.open(path)
    try:
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text_blocks.append(pytesseract.image_to_string(image))
    finally:
        doc.close()
    return "\n".join(text_blocks).strip()


def _extract_docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs).strip()


def _extract_image_text(path: Path) -> str:
    import pytesseract
    from PIL import Image

    image = Image.open(path)
    return pytesseract.image_to_string(image).strip()
